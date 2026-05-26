"""
PVM MineStar Weekly Performance Report Generator
Capstone Copper — built for engineer self-service
"""

import streamlit as st
import pandas as pd
import io
import os
import base64
from datetime import date, timedelta
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand colours ───────────────────────────────────────────────────────────────
CC_NAVY   = RGBColor(0x05, 0x2B, 0x48)
CC_ORANGE = RGBColor(0xD0, 0x5F, 0x27)
CC_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CC_LGRAY  = RGBColor(0xE7, 0xE6, 0xE6)
CC_MGRAY  = RGBColor(0x54, 0x58, 0x60)
CC_DTEXT  = RGBColor(0x1E, 0x29, 0x3B)
CC_RED    = RGBColor(0xC0, 0x39, 0x2B)
CC_GREEN  = RGBColor(0x16, 0xA3, 0x4A)
CC_AMBER  = RGBColor(0xD9, 0x77, 0x06)

LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "capstone_logo.png")

# ── Page config ─────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="MineStar Weekly Report Generator",
    page_icon="⛏️",
    layout="wide",
    initial_sidebar_state="collapsed",
)

st.markdown("""
<style>
  .block-container { padding-top: 1.5rem; padding-bottom: 2rem; }
  h1, h2, h3 { font-family: Arial, sans-serif; }

  .brand-bar {
    background: #052B48; color: white;
    padding: 14px 24px; border-radius: 8px;
    display: flex; align-items: center; justify-content: space-between;
    margin-bottom: 1.5rem;
  }
  .brand-bar h1 { color: white; margin: 0; font-size: 1.4rem; }
  .brand-bar span { color: #D05F27; font-weight: bold; font-size: 0.9rem; letter-spacing: 1px; }

  .step-card {
    border: 1.5px solid #e0e0e0; border-radius: 10px;
    padding: 1.2rem 1.4rem; margin-bottom: 1rem;
    background: white;
  }
  .step-header {
    display: flex; align-items: center; gap: 10px;
    margin-bottom: 0.6rem;
  }
  .step-num {
    background: #052B48; color: white;
    width: 28px; height: 28px; border-radius: 50%;
    display: flex; align-items: center; justify-content: center;
    font-weight: bold; font-size: 0.85rem; flex-shrink: 0;
  }
  .step-num.done { background: #16A34A; }
  .step-title { font-weight: bold; font-size: 1rem; color: #052B48; }

  .instr-box {
    background: #f0f4f8; border-left: 4px solid #D05F27;
    padding: 10px 14px; border-radius: 0 6px 6px 0;
    margin: 8px 0 12px 0; font-size: 0.88rem; color: #333;
    line-height: 1.6;
  }
  .instr-box b { color: #052B48; }
  .date-chip {
    display: inline-block; background: #D05F27; color: white;
    padding: 2px 10px; border-radius: 12px; font-size: 0.8rem;
    font-weight: bold; margin: 4px 0;
  }

  div[data-testid="stButton"] > button {
    background: #D05F27; color: white; font-weight: bold;
    border: none; padding: 0.6rem 2.4rem; border-radius: 8px;
    font-size: 1.05rem; width: 100%;
  }
  div[data-testid="stButton"] > button:hover { background: #b84d1f; }

  .success-banner {
    background: #052B48; color: white; padding: 16px 24px;
    border-radius: 8px; text-align: center; margin: 1rem 0;
  }
  .success-banner h3 { color: #D05F27; margin: 0 0 4px 0; }
</style>
""", unsafe_allow_html=True)


# ── Helpers ─────────────────────────────────────────────────────────────────────
def week_range(report_date: date):
    mon = report_date - timedelta(days=report_date.weekday())
    sun = mon + timedelta(days=6)
    return mon, sun

def fmt_date(d: date) -> str:
    return d.strftime("%-d %b %Y")

def parse_pct(val):
    try:
        return float(str(val).replace("%", "").strip())
    except Exception:
        return None

def fmt_num(val, decimals=1):
    """Format a numeric value to `decimals` decimal places; return raw string if not numeric."""
    try:
        raw = str(val).replace(",", "")
        is_pct = "%" in raw
        n = float(raw.replace("%", ""))
        return f"{n:.{decimals}f}%" if is_pct else f"{n:.{decimals}f}"
    except Exception:
        return str(val)

def detect_loader_cols(df: pd.DataFrame):
    cols_lower = {c.lower(): c for c in df.columns}
    patterns = {
        "machine_class": ["loading class", "primary machine class", "machine class"],
        "avail":         ["avail %", "availability %", "avail%"],
        "op_eff":        ["op eff %", "op efficency %", "op efficiency %"],
        "uoa":           ["uoa %", "uoa%"],
        "loads":         ["loads", "completed cycle", "sum(mssumm.completed cycle)"],
        "ttm":           ["ttm (dtm)", "ttm (dmt)", "sum(mssumm.paylod (dmt))", "paylod"],
        "hang":          ["hang", "avg hang", "hang time", "avg(mssumm.hang time)"],
        "pr":            ["pr", "productive rate", "productice rate"],
        "good_avail":    ["good avail"],
        "good_uoa":      ["good uoa"],
        "good_oe":       ["good oe"],
        "good_hang":     ["good hang"],
    }
    mapping = {}
    for key, candidates in patterns.items():
        for cand in candidates:
            if cand in cols_lower:
                mapping[key] = cols_lower[cand]
                break
    return mapping

def detect_delay_cols(df: pd.DataFrame):
    cols_lower = {c.lower(): c for c in df.columns}
    patterns = {
        "machine_class": ["primary machine class", "machine class"],
        "machine":       ["primary machine", "machine"],
        "delay":         ["delay"],
        "duration_min":  ["sum of shift duration (min)", "shift duration (min)", "duration (min)"],
    }
    mapping = {}
    for key, candidates in patterns.items():
        for cand in candidates:
            if cand in cols_lower:
                mapping[key] = cols_lower[cand]
                break
    return mapping

def detect_drill_cols(df: pd.DataFrame):
    cols_lower = {c.lower(): c for c in df.columns}
    patterns = {
        "machine_class": ["primary machine class", "loading class", "machine class"],
        "avail":         ["avail %", "availability %", "avail%"],
        "op_eff":        ["op eff %", "op efficency %", "op efficiency %"],
        "uoa":           ["uoa %", "uoa%"],
        "metres":        ["metres drilled", "meters drilled", "total metres", "total meters",
                          "metres", "meters", "drilled (m)"],
        "good_avail":    ["good avail"],
        "good_oe":       ["good oe"],
        "good_uoa":      ["good uoa"],
    }
    mapping = {}
    for key, candidates in patterns.items():
        for cand in candidates:
            if cand in cols_lower:
                mapping[key] = cols_lower[cand]
                break
    return mapping


# ── PPTX helpers ────────────────────────────────────────────────────────────────
def pptx_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def pptx_rect(slide, x, y, w, h, fill_rgb, line=False):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    if not line:
        s.line.fill.background()
    return s

def pptx_txt(slide, text, x, y, w, h, pt, rgb, bold=False,
             align=PP_ALIGN.LEFT, italic=False, wrap=True, margin=0.0, font="Arial"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size = Pt(pt); r.font.color.rgb = rgb
    r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tb

def pptx_logo(slide, x=8.55, y=0.06, w=1.25, h=0.42):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(x), Inches(y), Inches(w), Inches(h))

def pptx_header(slide, title, subtitle=""):
    pptx_rect(slide, 0, 0, 10, 0.64, CC_NAVY)
    pptx_rect(slide, 0, 0.64, 10, 0.04, CC_ORANGE)
    pptx_txt(slide, title, 0.25, 0.08, 7.5, 0.5, 16, CC_WHITE, bold=True, margin=0)
    if subtitle:
        pptx_txt(slide, subtitle, 0.25, 0.41, 7, 0.22, 8.5,
                 RGBColor(0xC0, 0xCC, 0xD8), margin=0)
    pptx_logo(slide)

def pptx_footer(slide, week_str):
    pptx_rect(slide, 0, 5.42, 10, 0.205, CC_NAVY)
    pptx_txt(slide, f"PVM MineStar  |  {week_str}  |  Capstone Copper  |  Confidential",
             0.3, 5.44, 9.4, 0.18, 7.5, RGBColor(0xAA, 0xBB, 0xCC), margin=0)

def kpi_card(slide, x, y, w, h, label, value, accent=None, value_pt=26, label_pt=8.5):
    accent = accent or CC_ORANGE
    pptx_rect(slide, x, y, w, h, CC_WHITE)
    pptx_rect(slide, x, y, w, 0.055, accent)
    pptx_txt(slide, value, x + 0.12, y + 0.1, w - 0.2, h * 0.58,
             value_pt, CC_NAVY, bold=True)
    pptx_txt(slide, label, x + 0.12, y + h * 0.62, w - 0.2, h * 0.4,
             label_pt, CC_MGRAY)


# ── Slide builders ──────────────────────────────────────────────────────────────
def build_title_slide(prs, week_str, report_date):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = CC_NAVY
    pptx_rect(s, 0, 0, 0.18, 5.625, CC_ORANGE)
    pptx_rect(s, 0, 5.3, 10, 0.325, RGBColor(0x03, 0x1A, 0x2C))
    pptx_txt(s, "PVM MINESTAR", 0.5, 0.8, 7, 0.4, 11, CC_ORANGE, bold=True, margin=0)
    pptx_txt(s, "Weekly Performance Report", 0.5, 1.22, 8, 0.9, 34,
             CC_WHITE, bold=True, margin=0)
    pptx_txt(s, week_str, 0.5, 2.18, 6, 0.35, 13, RGBColor(0xC0, 0xCC, 0xD8), margin=0)
    pptx_txt(s, f"Generated {report_date.strftime('%d %b %Y')}", 0.5, 2.56, 5, 0.28, 10,
             RGBColor(0x80, 0x90, 0xA0), margin=0)
    pptx_logo(s, x=0.5, y=4.5, w=2.0, h=0.68)
    pptx_txt(s, "CAPSTONE COPPER CONFIDENTIAL", 0.5, 5.34, 6, 0.2, 7.5,
             RGBColor(0x55, 0x66, 0x77), margin=0)


def build_executive_summary(prs, week_str, loader_df, loader_map, delay_df, delay_map, bad_df):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "EXECUTIVE SUMMARY", week_str)
    pptx_footer(s, week_str)

    cards = []
    if loader_df is not None and loader_map:
        ttm_col  = loader_map.get("ttm")
        loads_col = loader_map.get("loads")
        hang_col = loader_map.get("hang")
        if ttm_col:
            total_ttm = loader_df[ttm_col].apply(
                lambda x: float(str(x).replace(",", "")) if pd.notna(x) else 0).sum()
            cards.append(("Total TTM (dmt)", f"{total_ttm:,.0f}", CC_NAVY))
        if loads_col:
            total_loads = loader_df[loads_col].apply(
                lambda x: float(str(x).replace(",", "")) if pd.notna(x) else 0).sum()
            cards.append(("Total Loader Passes", f"{total_loads:,.0f}", CC_NAVY))
        if hang_col:
            avg_hang = loader_df[hang_col].apply(
                lambda x: float(str(x).replace(",", "")) if pd.notna(x) else None).mean()
            if avg_hang is not None:
                cards.append(("Avg Hang Time (min)", f"{avg_hang:.1f}", CC_AMBER))
    alert_count = len(bad_df) if bad_df is not None else 0
    cards.append(("Alarm Events", str(alert_count), CC_RED if alert_count > 0 else CC_GREEN))

    card_w = 9.4 / max(len(cards), 1)
    for i, (lbl, val, col) in enumerate(cards):
        kpi_card(s, 0.3 + i * (card_w + 0.05), 0.82, card_w - 0.05, 0.82, lbl, val, col)

    findings = []
    if loader_df is not None and loader_map:
        avail_col = loader_map.get("avail")
        if avail_col:
            for _, row in loader_df.iterrows():
                mc = row.get(loader_map.get("machine_class", ""), "")
                av = parse_pct(row.get(avail_col, ""))
                if av is not None and av < 80:
                    findings.append(("CRITICAL", CC_RED,
                        f"{mc} availability at {av:.1f}% — critically below target."))
        uoa_col = loader_map.get("uoa")
        if uoa_col:
            below_uoa = []
            for _, row in loader_df.iterrows():
                mc = row.get(loader_map.get("machine_class", ""), "")
                gd = str(row.get(loader_map.get("good_uoa", ""), "")).upper()
                uv = parse_pct(row.get(uoa_col, ""))
                if "D645" in gd or "FF0000" in gd or (uv is not None and uv < 70):
                    below_uoa.append(mc)
            if below_uoa:
                findings.append(("WARNING", CC_AMBER,
                    f"UoA below target: {', '.join(below_uoa)}."))

    if delay_df is not None and delay_map:
        dc = delay_map.get("machine_class")
        dd = delay_map.get("delay")
        dm = delay_map.get("duration_min")
        if dc and dd and dm:
            delay_df[dm] = pd.to_numeric(delay_df[dm], errors="coerce").fillna(0)
            no_op = delay_df[delay_df[dd].str.contains("No Operator", na=False)].groupby(dc)[dm].sum()
            for mclass, mins in no_op.sort_values(ascending=False).head(1).items():
                if mins > 300:
                    findings.append(("CRITICAL", CC_RED,
                        f"No Operator: {mclass} — {mins/60:.1f} hrs. Manning gap requires action."))
            blocked = delay_df[delay_df[dd].str.contains("Blocked Access", na=False)].groupby(dc)[dm].sum()
            for mclass, mins in blocked.sort_values(ascending=False).head(1).items():
                if mins > 200:
                    findings.append(("WARNING", CC_AMBER,
                        f"{mclass} Blocked Access: {mins/60:.1f} hrs — co-ordinate blast & road scheduling."))

    if bad_df is not None and len(bad_df) > 0:
        findings.append(("WARNING", CC_AMBER,
            f"{len(bad_df)} alarm event(s) recorded this week. See Bad Practices slide."))
    if not findings:
        findings.append(("INFO", CC_NAVY, "No critical issues detected this week."))

    fy = 1.76
    for sev, col, msg in findings[:6]:
        pptx_rect(s, 0.3, fy, 9.4, 0.62, RGBColor(0xEE, 0xF2, 0xF8))
        pptx_rect(s, 0.3, fy, 0.07, 0.62, col)
        pptx_txt(s, sev, 0.45, fy + 0.04, 1.2, 0.2, 7.5, col, bold=True, margin=0)
        pptx_txt(s, msg, 0.45, fy + 0.28, 9.1, 0.32, 9, CC_DTEXT, wrap=True, margin=0)
        fy += 0.70


def _draw_kpi_table(s, df, display_cols, first_col_key, th_y=1.72, rh=0.46, font_size=9.5):
    """Generic table renderer used by loader/truck/drill slides."""
    total_w = sum(w for _, _, w in display_cols)
    scale   = 9.4 / total_w
    scaled  = [(col, hdr, w * scale) for col, hdr, w in display_cols]
    xs = [0.3]
    for _, _, w in scaled[:-1]:
        xs.append(xs[-1] + w)

    pptx_rect(s, 0.3, th_y, 9.4, rh * 0.88, CC_NAVY)
    for (col, hdr, w), x in zip(scaled, xs):
        align = PP_ALIGN.LEFT if col == first_col_key else PP_ALIGN.CENTER
        pptx_txt(s, hdr, x + 0.06, th_y + 0.1, w - 0.08, 0.28,
                 8.5, CC_WHITE, bold=True, align=align, margin=0)
    return scaled, xs


def build_loader_kpi_slide(prs, week_str, loader_df, loader_map):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "LOADER KPI PERFORMANCE", week_str)
    pptx_footer(s, week_str)

    if loader_df is None:
        pptx_txt(s, "No loader data uploaded.", 0.5, 2.5, 9, 0.5, 14, CC_MGRAY)
        return

    mc_col  = loader_map.get("machine_class", loader_df.columns[0])
    avail_c = loader_map.get("avail")
    opeff_c = loader_map.get("op_eff")
    uoa_c   = loader_map.get("uoa")
    loads_c = loader_map.get("loads")
    ttm_c   = loader_map.get("ttm")
    hang_c  = loader_map.get("hang")
    g_avail = loader_map.get("good_avail")
    g_uoa   = loader_map.get("good_uoa")
    g_hang  = loader_map.get("good_hang")

    ttm_total = loader_df[ttm_c].apply(
        lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum() if ttm_c else 0
    loads_total = loader_df[loads_c].apply(
        lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum() if loads_c else 0
    hang_vals = loader_df[hang_c].apply(
        lambda x: float(str(x).replace(",","")) if pd.notna(x) else None).dropna() if hang_c else pd.Series([])
    hang_avg = hang_vals.mean() if len(hang_vals) else None

    cards = [
        ("Total TTM (dmt)", f"{ttm_total:,.0f}", CC_NAVY),
        ("Total Loader Passes", f"{loads_total:,.0f}", CC_NAVY),
        ("Avg Hang Time (min)", f"{hang_avg:.1f}" if hang_avg is not None else "—", CC_AMBER),
        ("Loaders Reporting", str(len(loader_df)), CC_NAVY),
    ]
    cw = 2.2
    for i, (lbl, val, col) in enumerate(cards):
        kpi_card(s, 0.3 + i * 2.35, 0.82, cw, 0.75, lbl, val, col, value_pt=24)

    display_cols = [(mc_col, "Machine Class", 2.0)]
    if avail_c:  display_cols.append((avail_c,  "Avail %",    0.85))
    if opeff_c:  display_cols.append((opeff_c,  "Op Eff %",   0.9))
    if uoa_c:    display_cols.append((uoa_c,    "UoA %",      0.85))
    if loads_c:  display_cols.append((loads_c,  "Loads",      0.75))
    if ttm_c:    display_cols.append((ttm_c,    "TTM (dmt)",  1.4))
    if hang_c:   display_cols.append((hang_c,   "Hang (min)", 0.95))

    th_y, rh = 1.72, 0.46
    scaled, xs = _draw_kpi_table(s, loader_df, display_cols, mc_col, th_y, rh)

    for ri, (_, row) in enumerate(loader_df.iterrows()):
        ry = th_y + rh * (ri + 1)
        pptx_rect(s, 0.3, ry, 9.4, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
        for (col, hdr, w), x in zip(scaled, xs):
            raw = str(row.get(col, "—"))
            cell_col = CC_DTEXT
            cell_bold = (col == mc_col)
            align = PP_ALIGN.LEFT if col == mc_col else PP_ALIGN.CENTER

            def _gv(gcol):
                return str(row.get(gcol, "")).upper() if gcol and gcol in row.index else ""

            if col == avail_c:
                gv = _gv(g_avail)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)
            elif col == uoa_c:
                gv = _gv(g_uoa)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)
            elif col == hang_c:
                gv = _gv(g_hang)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)

            # Format values
            if col == ttm_c:
                try: raw = f"{float(raw.replace(',','')):.0f}"
                except Exception: pass
            elif col != mc_col:
                raw = fmt_num(raw)

            pptx_txt(s, raw, x + 0.06, ry + 0.1, w - 0.08, 0.28,
                     9.5, cell_col, bold=cell_bold, align=align, margin=0)

    leg_y = th_y + rh * (len(loader_df) + 1) + 0.08
    pptx_rect(s, 0.3, leg_y, 0.16, 0.16, CC_GREEN)
    pptx_txt(s, "= At or above target", 0.5, leg_y - 0.01, 2.5, 0.2, 7.5, CC_MGRAY, margin=0)
    pptx_rect(s, 2.9, leg_y, 0.16, 0.16, CC_RED)
    pptx_txt(s, "= Below target", 3.1, leg_y - 0.01, 2.5, 0.2, 7.5, CC_MGRAY, margin=0)


def build_truck_kpi_slide(prs, week_str, truck_df, truck_map):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "TRUCK KPI PERFORMANCE", week_str)
    pptx_footer(s, week_str)

    if truck_df is None:
        pptx_txt(s, "No truck data uploaded.", 0.5, 2.5, 9, 0.5, 14, CC_MGRAY)
        return

    mc_col = truck_map.get("machine_class", truck_df.columns[0])
    avail_c = truck_map.get("avail")
    opeff_c = truck_map.get("op_eff")
    uoa_c   = truck_map.get("uoa")
    loads_c = truck_map.get("loads")
    ttm_c   = truck_map.get("ttm")

    # Use "Truck Class" column as first column; hide raw Machine Class
    import re as _re
    truck_class_col = next(
        (c for c in truck_df.columns if _re.search(r'truck.?class|class.?truck', c, _re.I)),
        mc_col
    )

    cards = []
    if ttm_c:
        ttm_total = truck_df[ttm_c].apply(lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum()
        cards.append(("Total Payload (dmt)", f"{ttm_total:,.0f}", CC_NAVY))
    if loads_c:
        loads_total = truck_df[loads_c].apply(lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum()
        cards.append(("Total Cycles", f"{loads_total:,.0f}", CC_NAVY))
    cards.append(("Truck Classes", str(len(truck_df)), CC_NAVY))

    cw = 9.4 / max(len(cards), 1) - 0.1
    for i, (lbl, val, col) in enumerate(cards):
        kpi_card(s, 0.3 + i * (cw + 0.1), 0.82, cw, 0.75, lbl, val, col, value_pt=24)

    # Truck Class first; Machine Class excluded
    display_cols = [(truck_class_col, "Truck Class", 2.0)]
    if avail_c:  display_cols.append((avail_c, "Avail %",  0.85))
    if opeff_c:  display_cols.append((opeff_c, "Op Eff %", 0.9))
    if uoa_c:    display_cols.append((uoa_c,   "UoA %",    0.85))
    if loads_c:  display_cols.append((loads_c, "Cycles",   0.85))
    if ttm_c:    display_cols.append((ttm_c,   "Payload (dmt)", 1.4))
    used = {truck_class_col, mc_col, avail_c, opeff_c, uoa_c, loads_c, ttm_c}
    for col in truck_df.columns:
        if col not in used and "good" not in col.lower() and len(display_cols) < 8:
            display_cols.append((col, col[:14], 1.0))

    th_y, rh = 1.72, 0.46
    scaled, xs = _draw_kpi_table(s, truck_df, display_cols, truck_class_col, th_y, rh)

    for ri, (_, row) in enumerate(truck_df.iterrows()):
        ry = th_y + rh * (ri + 1)
        pptx_rect(s, 0.3, ry, 9.4, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
        for (col, hdr, w), x in zip(scaled, xs):
            raw = str(row.get(col, "—"))
            cell_col = CC_DTEXT
            cell_bold = (col == truck_class_col)
            align = PP_ALIGN.LEFT if col == truck_class_col else PP_ALIGN.CENTER
            if col == avail_c:
                pv = parse_pct(raw)
                if pv is not None: cell_col = CC_GREEN if pv >= 85 else CC_RED
            elif col == uoa_c:
                pv = parse_pct(raw)
                if pv is not None: cell_col = CC_GREEN if pv >= 70 else CC_RED
            if col != truck_class_col:
                raw = fmt_num(raw)
            pptx_txt(s, raw, x + 0.06, ry + 0.1, w - 0.08, 0.28,
                     9, cell_col, bold=cell_bold, align=align, margin=0)


def build_drill_kpi_slide(prs, week_str, drill_df, drill_map):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "DRILL KPI PERFORMANCE", week_str)
    pptx_footer(s, week_str)

    if drill_df is None:
        pptx_txt(s, "No drill data uploaded.", 0.5, 2.5, 9, 0.5, 14, CC_MGRAY)
        return

    mc_col  = drill_map.get("machine_class", drill_df.columns[0])
    avail_c = drill_map.get("avail")
    opeff_c = drill_map.get("op_eff")
    uoa_c   = drill_map.get("uoa")
    mtr_c   = drill_map.get("metres")
    g_avail = drill_map.get("good_avail")
    g_oe    = drill_map.get("good_oe")
    g_uoa   = drill_map.get("good_uoa")

    metres_total = drill_df[mtr_c].apply(
        lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum() if mtr_c else 0
    oe_vals = drill_df[opeff_c].apply(parse_pct).dropna() if opeff_c else pd.Series([])
    avg_oe  = oe_vals.mean() if len(oe_vals) else None

    cards = []
    if metres_total > 0:
        cards.append(("Total Metres Drilled", f"{metres_total:,.1f}", CC_NAVY))
    if avg_oe is not None:
        cards.append(("Avg Op Eff %", f"{avg_oe:.1f}%", CC_AMBER))
    cards.append(("Drill Classes", str(len(drill_df)), CC_NAVY))

    cw = 9.4 / max(len(cards), 1) - 0.1
    for i, (lbl, val, col) in enumerate(cards):
        kpi_card(s, 0.3 + i * (cw + 0.1), 0.82, cw, 0.75, lbl, val, col, value_pt=24)

    display_cols = [(mc_col, "Machine Class", 2.0)]
    if avail_c:  display_cols.append((avail_c, "Avail %",  0.85))
    if opeff_c:  display_cols.append((opeff_c, "Op Eff %", 0.9))
    if uoa_c:    display_cols.append((uoa_c,   "UoA %",    0.85))
    if mtr_c:    display_cols.append((mtr_c,   "Metres",   1.2))
    used = {mc_col, avail_c, opeff_c, uoa_c, mtr_c}
    for col in drill_df.columns:
        if col not in used and "good" not in col.lower() and len(display_cols) < 7:
            display_cols.append((col, col[:14], 1.0))

    th_y, rh = 1.72, 0.46
    scaled, xs = _draw_kpi_table(s, drill_df, display_cols, mc_col, th_y, rh)

    for ri, (_, row) in enumerate(drill_df.iterrows()):
        ry = th_y + rh * (ri + 1)
        pptx_rect(s, 0.3, ry, 9.4, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
        for (col, hdr, w), x in zip(scaled, xs):
            raw = str(row.get(col, "—"))
            cell_col = CC_DTEXT
            cell_bold = (col == mc_col)
            align = PP_ALIGN.LEFT if col == mc_col else PP_ALIGN.CENTER

            def _gv(gcol):
                return str(row.get(gcol, "")).upper() if gcol and gcol in row.index else ""

            if col == avail_c:
                gv = _gv(g_avail)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)
            elif col == uoa_c:
                gv = _gv(g_uoa)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)
            elif col == opeff_c:
                gv = _gv(g_oe)
                cell_col = CC_GREEN if "00FF00" in gv else (CC_RED if ("D645" in gv or "FF0000" in gv) else CC_DTEXT)

            if col != mc_col:
                raw = fmt_num(raw)

            pptx_txt(s, raw, x + 0.06, ry + 0.1, w - 0.08, 0.28,
                     9, cell_col, bold=cell_bold, align=align, margin=0)

    leg_y = th_y + rh * (len(drill_df) + 1) + 0.08
    pptx_rect(s, 0.3, leg_y, 0.16, 0.16, CC_GREEN)
    pptx_txt(s, "= At or above target", 0.5, leg_y - 0.01, 2.5, 0.2, 7.5, CC_MGRAY, margin=0)
    pptx_rect(s, 2.9, leg_y, 0.16, 0.16, CC_RED)
    pptx_txt(s, "= Below target", 3.1, leg_y - 0.01, 2.5, 0.2, 7.5, CC_MGRAY, margin=0)


def build_efficiency_slide(prs, week_str, loader_df, loader_map,
                            truck_df, truck_map, drill_df, drill_map,
                            delay_df, delay_map):
    """Efficiency derived from equipment OE columns + delay data — no extra upload needed."""
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "EFFICIENCY OVERVIEW", week_str)
    pptx_footer(s, week_str)

    # OE cards per equipment type
    equip_sets = [
        ("Loaders", loader_df, loader_map),
        ("Trucks",  truck_df,  truck_map),
        ("Drills",  drill_df,  drill_map),
    ]
    oe_cards = []
    for label, df, m in equip_sets:
        if df is None or not m: continue
        oe_col = m.get("op_eff")
        if not oe_col: continue
        vals = df[oe_col].apply(parse_pct).dropna()
        if len(vals) == 0: continue
        avg = vals.mean()
        oe_cards.append((f"{label} OE", f"{avg:.1f}%", CC_GREEN if avg >= 75 else CC_RED))

    if not oe_cards:
        oe_cards = [("No Equipment Data", "—", CC_MGRAY)]

    cw = 9.4 / max(len(oe_cards), 1) - 0.1
    for i, (lbl, val, col) in enumerate(oe_cards):
        kpi_card(s, 0.3 + i * (cw + 0.1), 0.82, cw, 0.75, lbl, val, col, value_pt=28)

    # Top delay types table
    pptx_txt(s, "Top Delay Types — Hours Lost by Equipment",
             0.3, 1.72, 9.4, 0.3, 10, CC_NAVY, bold=True, margin=0)

    if delay_df is not None and delay_map:
        dd = delay_map.get("delay")
        dm = delay_map.get("duration_min")
        if dd and dm:
            delay_df[dm] = pd.to_numeric(delay_df[dm], errors="coerce").fillna(0)
            agg = delay_df.groupby(dd)[dm].sum().sort_values(ascending=False).head(8)
            grand_total = agg.sum() or 1

            hdrs    = ["Delay Type", "Total Hours", "% of Total"]
            col_ws  = [5.5, 1.8, 1.8]
            th_y, rh = 2.08, 0.38
            xs = [0.3]; [xs.append(xs[-1] + w) for w in col_ws[:-1]]

            pptx_rect(s, 0.3, th_y, 9.1, rh * 0.88, CC_NAVY)
            for hdr, w, x in zip(hdrs, col_ws, xs):
                align = PP_ALIGN.LEFT if x == xs[0] else PP_ALIGN.CENTER
                pptx_txt(s, hdr, x + 0.05, th_y + 0.08, w - 0.06, 0.24,
                         8.5, CC_WHITE, bold=True, align=align, margin=0)

            for ri, (delay_name, mins) in enumerate(agg.items()):
                ry = th_y + rh * (ri + 1)
                pptx_rect(s, 0.3, ry, 9.1, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
                vals = [str(delay_name)[:35], f"{mins/60:.1f}", f"{mins/grand_total*100:.1f}%"]
                for val, w, x in zip(vals, col_ws, xs):
                    align = PP_ALIGN.LEFT if x == xs[0] else PP_ALIGN.CENTER
                    pptx_txt(s, val, x + 0.05, ry + 0.08, w - 0.06, 0.24,
                             8.5, CC_DTEXT, bold=(x == xs[0]), align=align, margin=0)
        else:
            pptx_txt(s, "Delay columns could not be detected.", 0.3, 2.1, 9.4, 0.4, 10, CC_RED)
    else:
        pptx_txt(s, "No delay data uploaded.", 0.3, 2.1, 9.4, 0.4, 10, CC_MGRAY)


def build_delay_slide(prs, week_str, delay_df, delay_map):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "DELAY ANALYSIS", week_str)
    pptx_footer(s, week_str)

    if delay_df is None:
        pptx_txt(s, "No delay data uploaded.", 0.5, 2.5, 9, 0.5, 14, CC_MGRAY)
        return

    dc = delay_map.get("machine_class")
    dd = delay_map.get("delay")
    dm = delay_map.get("duration_min")
    if not (dc and dd and dm):
        pptx_txt(s, "Delay data columns could not be detected.", 0.5, 2.5, 9, 0.5, 12, CC_RED)
        return

    delay_df[dm] = pd.to_numeric(delay_df[dm], errors="coerce").fillna(0)
    agg = defaultdict(lambda: defaultdict(float))
    for _, row in delay_df.iterrows():
        agg[row[dc]][row[dd]] += row[dm]

    rows = []
    for mclass, delays in sorted(agg.items(), key=lambda x: -sum(x[1].values())):
        total_hrs = sum(delays.values()) / 60
        sorted_d  = sorted(delays.items(), key=lambda x: -x[1])
        top1 = sorted_d[0] if sorted_d else ("—", 0)
        top2 = sorted_d[1] if len(sorted_d) > 1 else ("—", 0)
        sev  = CC_RED if total_hrs > 1000 else (CC_AMBER if total_hrs > 300 else CC_MGRAY)
        rows.append((mclass, f"{total_hrs:.1f}", top1[0][:28], f"{top1[1]/60:.1f}",
                     top2[0][:28], f"{top2[1]/60:.1f}", sev))

    alert_rows = [r for r in rows if r[6] in (CC_RED, CC_AMBER)][:4]
    card_w = 9.4 / max(len(alert_rows), 1) - 0.12
    for i, r in enumerate(alert_rows):
        cx = 0.3 + i * (card_w + 0.12)
        pptx_rect(s, cx, 0.82, card_w, 0.82, CC_WHITE)
        pptx_rect(s, cx, 0.82, card_w, 0.055, r[6])
        pptx_txt(s, f"{r[1]} hrs", cx + 0.1, 0.9, card_w - 0.15, 0.38, 22, r[6], bold=True, margin=0)
        pptx_txt(s, r[0], cx + 0.1, 1.28, card_w - 0.15, 0.2, 8, CC_DTEXT, bold=True, wrap=True, margin=0)
        pptx_txt(s, r[2], cx + 0.1, 1.48, card_w - 0.15, 0.16, 7.5, CC_MGRAY, wrap=True, margin=0)

    headers = ["Machine Class", "Total Hrs", "Top Delay #1", "Hrs", "Top Delay #2", "Hrs"]
    col_ws  = [2.0, 0.85, 2.5, 0.7, 2.5, 0.7]
    th_y, rh = 1.82, 0.4
    xs = [0.3]; [xs.append(xs[-1] + w) for w in col_ws[:-1]]

    pptx_rect(s, 0.3, th_y, 9.25, rh * 0.88, CC_NAVY)
    for hdr, w, x in zip(headers, col_ws, xs):
        align = PP_ALIGN.LEFT if x == xs[0] else PP_ALIGN.CENTER
        pptx_txt(s, hdr, x + 0.05, th_y + 0.09, w - 0.06, 0.24,
                 8, CC_WHITE, bold=True, align=align, margin=0)

    for ri, row in enumerate(rows[:8]):
        ry = th_y + rh * (ri + 1)
        pptx_rect(s, 0.3, ry, 9.25, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
        for val, w, x, idx in zip([row[0],row[1],row[2],row[3],row[4],row[5]], col_ws, xs, range(6)):
            col_c = CC_DTEXT
            if idx in (1, 3, 5): col_c = row[6]
            align = PP_ALIGN.LEFT if x == xs[0] else PP_ALIGN.CENTER
            pptx_txt(s, str(val)[:30], x + 0.05, ry + 0.09, w - 0.06, 0.24,
                     8.5, col_c, bold=(x == xs[0]), align=align, margin=0)


def build_bad_practices_slide(prs, week_str, bad_df):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    pptx_header(s, "BAD PRACTICES & ALARMS", week_str)
    pptx_footer(s, week_str)

    if bad_df is None or len(bad_df) == 0:
        pptx_txt(s, "No bad practices data uploaded.", 0.5, 2.5, 9, 0.5, 12, CC_MGRAY)
        return

    n_events   = len(bad_df)
    n_machines = bad_df["Machine"].nunique() if "Machine" in bad_df.columns else "—"
    n_ops      = bad_df["Operator"].nunique() if "Operator" in bad_df.columns else "—"
    dur_col    = next((c for c in bad_df.columns if "duration" in c.lower()), None)
    total_dur  = bad_df[dur_col].apply(
        lambda x: float(str(x).replace(",","")) if pd.notna(x) else 0).sum() if dur_col else 0

    cards = [
        ("Alarm Events",             str(n_events),        CC_RED),
        ("Machines Affected",        str(n_machines),      CC_AMBER),
        ("Operators Involved",       str(n_ops),           CC_AMBER),
        ("Total Alarm Duration (s)", f"{total_dur:.1f}",   CC_RED),
    ]
    cw = 2.2
    for i, (lbl, val, col) in enumerate(cards):
        kpi_card(s, 0.3 + i * 2.35, 0.82, cw, 0.75, lbl, val, col, value_pt=24)

    display_cols = [c for c in ["TimeStamp", "Machine", "Operator", "Alarm Desc", dur_col or "Duration"]
                    if c and c in bad_df.columns]
    if not display_cols:
        display_cols = list(bad_df.columns[:5])

    col_ws_map = {"TimeStamp": 1.6, "Machine": 0.75, "Operator": 1.5, "Alarm Desc": 3.8}
    col_ws = [col_ws_map.get(c, 1.2) for c in display_cols]
    scale  = 9.4 / sum(col_ws)
    col_ws = [w * scale for w in col_ws]

    th_y, rh = 1.72, 0.44
    xs = [0.3]; [xs.append(xs[-1] + w) for w in col_ws[:-1]]

    pptx_rect(s, 0.3, th_y, 9.4, rh * 0.88, CC_NAVY)
    for hdr, w, x in zip(display_cols, col_ws, xs):
        pptx_txt(s, hdr, x + 0.05, th_y + 0.1, w - 0.06, 0.26, 8.5, CC_WHITE, bold=True, margin=0)

    for ri, (_, row) in enumerate(bad_df.head(10).iterrows()):
        ry = th_y + rh * (ri + 1)
        pptx_rect(s, 0.3, ry, 9.4, rh * 0.88, CC_WHITE if ri % 2 == 0 else CC_LGRAY)
        for col, w, x in zip(display_cols, col_ws, xs):
            val = str(row.get(col, ""))[:40]
            cell_col = CC_RED if col == "Alarm Desc" else CC_DTEXT
            pptx_txt(s, val, x + 0.05, ry + 0.1, w - 0.06, 0.26,
                     8.5, cell_col, bold=(col == "Machine"), margin=0, wrap=True)

    if len(bad_df) > 10:
        pptx_txt(s, f"… and {len(bad_df)-10} more events.",
                 0.3, th_y + rh * 11 + 0.05, 9.4, 0.25, 8, CC_MGRAY, margin=0)


def build_actions_slide(prs, week_str, loader_df, loader_map, delay_df, delay_map, bad_df):
    s = pptx_blank(prs)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = CC_NAVY
    pptx_rect(s, 0, 0, 10, 0.04, CC_ORANGE)
    pptx_txt(s, "ACTIONS & ISSUES", 0.25, 0.12, 7.5, 0.38, 16, CC_WHITE, bold=True, margin=0)
    pptx_txt(s, f"Items requiring follow-up  |  {week_str}", 0.25, 0.46, 7, 0.24, 9,
             RGBColor(0xAA, 0xBB, 0xCC), margin=0)
    pptx_logo(s)
    pptx_rect(s, 0, 5.42, 10, 0.205, RGBColor(0x03, 0x15, 0x28))
    pptx_txt(s, f"PVM MineStar  |  {week_str}  |  Capstone Copper  |  Confidential",
             0.3, 5.44, 9.4, 0.18, 7.5, RGBColor(0x55, 0x66, 0x77), margin=0)

    actions = []
    if loader_df is not None and loader_map:
        avail_c = loader_map.get("avail")
        mc_col  = loader_map.get("machine_class", loader_df.columns[0])
        if avail_c:
            for _, row in loader_df.iterrows():
                mc = str(row.get(mc_col, ""))
                av = parse_pct(row.get(avail_c, ""))
                if av is not None and av < 80:
                    actions.append((CC_RED, "CRITICAL", f"{mc} Availability at {av:.1f}%",
                        "Availability critically low. Review maintenance schedule and mechanical faults. Escalate to Maintenance Planner.",
                        "Maintenance Planner"))

    if delay_df is not None and delay_map:
        dc = delay_map.get("machine_class")
        dd = delay_map.get("delay")
        dm = delay_map.get("duration_min")
        if dc and dd and dm:
            delay_df[dm] = pd.to_numeric(delay_df[dm], errors="coerce").fillna(0)
            no_op = delay_df[delay_df[dd].str.contains("No Operator", na=False)] \
                .groupby(dc)[dm].sum().sort_values(ascending=False)
            for mclass, mins in no_op.head(2).items():
                if mins > 300:
                    actions.append((CC_RED, "CRITICAL",
                        f"{mclass} — No Operator ({mins/60:.1f} hrs)",
                        f"Manning gap of {mins/60:.1f} hrs. Review roster coverage and callout protocols urgently.",
                        "HR / Roster Planning"))
            blocked = delay_df[delay_df[dd].str.contains("Blocked Access", na=False)] \
                .groupby(dc)[dm].sum().sort_values(ascending=False)
            for mclass, mins in blocked.head(1).items():
                if mins > 200:
                    actions.append((CC_AMBER, "WARNING",
                        f"{mclass} — Blocked Access ({mins/60:.1f} hrs)",
                        "Co-ordinate blast scheduling and haul road availability with Drill & Blast and Roads teams.",
                        "Drill & Blast / Roads"))

    if bad_df is not None and len(bad_df) > 0:
        alarm_col = next((c for c in bad_df.columns if "alarm" in c.lower()), None)
        if alarm_col:
            for alarm, grp in bad_df.groupby(alarm_col):
                if len(grp) >= 2:
                    op  = grp["Operator"].iloc[0] if "Operator" in grp.columns else "Unknown"
                    mch = grp["Machine"].iloc[0] if "Machine" in grp.columns else "Unknown"
                    actions.append((CC_AMBER, "WARNING",
                        f"Machine {mch} — {str(alarm)[:40]}",
                        f"{len(grp)} repeated alarms. Operator: {op}. Inspect machine and review alarm response protocol.",
                        "Maintenance / Supervisor"))

    if not actions:
        actions.append((CC_NAVY, "INFO", "No critical actions this week",
            "All key metrics within acceptable range.", "Operations"))

    ay = 0.82
    for sev_col, sev_lbl, title, detail, owner in actions[:6]:
        pptx_rect(s, 0.3, ay, 9.4, 0.77, RGBColor(0x0D, 0x22, 0x42))
        pptx_rect(s, 0.3, ay, 0.07, 0.77, sev_col)
        pptx_rect(s, 0.45, ay + 0.08, 0.9, 0.22, sev_col)
        pptx_txt(s, sev_lbl, 0.46, ay + 0.09, 0.88, 0.2, 7, CC_WHITE, bold=True,
                 align=PP_ALIGN.CENTER, margin=0)
        pptx_txt(s, title, 1.45, ay + 0.06, 5.8, 0.26, 10, CC_WHITE, bold=True, margin=0)
        pptx_txt(s, f"Owner: {owner}", 7.3, ay + 0.06, 2.3, 0.22, 8, CC_ORANGE,
                 align=PP_ALIGN.RIGHT, margin=0)
        pptx_txt(s, detail, 1.45, ay + 0.36, 8.1, 0.38, 8.5,
                 RGBColor(0xBB, 0xCC, 0xDD), wrap=True, margin=0)
        ay += 0.83


# ── Main PPTX builder ───────────────────────────────────────────────────────────
def generate_report(loader_df, loader_map, truck_df, truck_map,
                    drill_df, drill_map, delay_df, delay_map,
                    bad_df, week_str, report_date):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)

    build_title_slide(prs, week_str, report_date)
    build_executive_summary(prs, week_str, loader_df, loader_map, delay_df, delay_map, bad_df)
    build_loader_kpi_slide(prs, week_str, loader_df, loader_map)
    build_truck_kpi_slide(prs, week_str, truck_df, truck_map)
    build_drill_kpi_slide(prs, week_str, drill_df, drill_map)
    build_efficiency_slide(prs, week_str, loader_df, loader_map,
                           truck_df, truck_map, drill_df, drill_map,
                           delay_df, delay_map)
    build_delay_slide(prs, week_str, delay_df, delay_map)
    build_bad_practices_slide(prs, week_str, bad_df)
    build_actions_slide(prs, week_str, loader_df, loader_map, delay_df, delay_map, bad_df)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── UI helpers ──────────────────────────────────────────────────────────────────
def upload_step(step_num, title, page_name, visual_name, columns_hint,
                key, week_mon, week_sun, required=True):
    uploaded = st.session_state.get(f"df_{key}")
    done = uploaded is not None
    num_cls = "done" if done else ""
    badge = "✅ Uploaded" if done else ("⬜ Required" if required else "⬜ Optional")
    badge_col = "#16A34A" if done else ("#D05F27" if required else "#888")

    st.markdown(f"""
    <div class="step-card">
      <div class="step-header">
        <div class="step-num {num_cls}">{step_num}</div>
        <div class="step-title">{title}</div>
        <span style="margin-left:auto;color:{badge_col};font-weight:bold;font-size:0.85rem">{badge}</span>
      </div>
      <div class="instr-box">
        <b>Power BI page:</b> <code>{page_name}</code> &nbsp;|&nbsp;
        <b>Visual:</b> <i>{visual_name}</i><br>
        <b>Date range:</b>
        <span class="date-chip">Mon {fmt_date(week_mon)} – Sun {fmt_date(week_sun)}</span><br><br>
        <b>How to export:</b>
        Click the visual → click <b>⋯ (More options)</b> → <b>Export data</b> →
        set format to <b>.csv</b> → click <b>Export</b>.<br>
        <b>Expected columns:</b> <code>{columns_hint}</code>
      </div>
    </div>
    """, unsafe_allow_html=True)

    files = st.file_uploader(f"Upload CSV(s) for: {title}", type=["csv"], key=f"up_{key}",
                             label_visibility="collapsed", accept_multiple_files=True)
    if files:
        try:
            dfs = [pd.read_csv(f, encoding="utf-8-sig") for f in files]
            df = pd.concat(dfs, ignore_index=True) if len(dfs) > 1 else dfs[0]
            st.session_state[f"df_{key}"] = df
            if len(files) > 1:
                st.success(f"✅  Merged {len(files)} files — {len(df)} total rows, {len(df.columns)} columns")
            else:
                st.success(f"✅  Loaded {len(df)} rows, {len(df.columns)} columns")
            with st.expander("Preview data"):
                st.dataframe(df.head(8), use_container_width=True)
            return df
        except Exception as e:
            st.error(f"Could not parse file(s): {e}")
    elif done:
        with st.expander("Preview uploaded data"):
            st.dataframe(uploaded.head(8), use_container_width=True)
        return uploaded
    return None


# ── Main app ────────────────────────────────────────────────────────────────────
def main():
    logo_b64 = ""
    if os.path.exists(LOGO_PATH):
        with open(LOGO_PATH, "rb") as f:
            logo_b64 = base64.b64encode(f.read()).decode()

    logo_html = f'<img src="data:image/png;base64,{logo_b64}" style="height:44px">' if logo_b64 else ""
    st.markdown(f"""
    <div class="brand-bar">
      <div><span>PVM MINESTAR</span><h1>Weekly Performance Report Generator</h1></div>
      {logo_html}
    </div>
    """, unsafe_allow_html=True)

    st.markdown("Upload your weekly Power BI CSV exports below and click **Generate Report** "
                "to produce a Capstone Copper-branded PowerPoint in seconds.")

    st.markdown("---")
    col_date, col_info = st.columns([2, 4])
    with col_date:
        st.markdown("**Select the Sunday of your reporting week:**")
        today = date.today()
        default_sunday = today + timedelta(days=(6 - today.weekday()) % 7)
        report_date = st.date_input("Week ending (Sunday)", value=default_sunday,
                                     label_visibility="collapsed")
    week_mon, week_sun = week_range(report_date)
    week_str = f"Week {fmt_date(week_mon)} – {fmt_date(week_sun)}"
    with col_info:
        st.info(f"📅  **Reporting week:** {week_str}\n\n"
                f"Filter Power BI exports to **Mon {fmt_date(week_mon)} – Sun {fmt_date(week_sun)}**.")

    st.markdown("---")
    st.markdown("### 📂 Upload Data Exports")

    loader_df = upload_step(
        1, "Loader Numbers",
        page_name="'Date Ranges'",
        visual_name="Loader summary table (Avail %, Op Eff %, UoA %, TTM, Loads, Hang)",
        columns_hint="Loading Class, Avail %, Op Eff %, UoA %, Loads, TTM (dtm), Hang",
        key="loader", week_mon=week_mon, week_sun=week_sun,
    )
    truck_df = upload_step(
        2, "Truck Numbers",
        page_name="'Date Ranges'",
        visual_name="Truck summary table (Avail %, UoA %, Op Eff %, Cycles, Payload)",
        columns_hint="Truck Class, Avail %, UoA %, Op Eff %, Completed Cycle, Payload (dmt)",
        key="truck", week_mon=week_mon, week_sun=week_sun,
    )
    drill_df = upload_step(
        3, "Drill Numbers",
        page_name="'Date Ranges'",
        visual_name="Drill summary table (Avail %, Op Eff %, UoA %, Metres Drilled)",
        columns_hint="Primary Machine Class, Avail %, Op Eff %, UoA %, Metres Drilled",
        key="drill", week_mon=week_mon, week_sun=week_sun,
    )
    delay_df = upload_step(
        4, "Delays Numbers",
        page_name="'Delays'",
        visual_name="Delay summary table (all machine classes & delay types)",
        columns_hint="Primary Machine Class, Primary Machine, Delay, Sum of Shift Duration (min)",
        key="delay", week_mon=week_mon, week_sun=week_sun,
    )
    bad_df = upload_step(
        5, "Bad Practices",
        page_name="'Bad Practices'",
        visual_name="Any bad practices table (Loader, Truck 789, Truck 793, Drill)",
        columns_hint="TimeStamp, Machine, Operator, Alarm Desc, Duration",
        key="bad", week_mon=week_mon, week_sun=week_sun,
    )

    st.markdown("---")
    required_dfs = [loader_df, truck_df, drill_df, delay_df, bad_df]
    required_names = ["Loader Numbers", "Truck Numbers", "Drill Numbers", "Delays Numbers", "Bad Practices"]
    missing = [n for n, df in zip(required_names, required_dfs) if df is None]

    if missing:
        st.warning(f"⚠️  Still waiting on: **{', '.join(missing)}**")

    if st.button("⚡  Generate PowerPoint Report", disabled=bool(missing)):
        with st.spinner("Building your Capstone Copper report…"):
            loader_map = detect_loader_cols(loader_df) if loader_df is not None else {}
            truck_map  = detect_loader_cols(truck_df)  if truck_df  is not None else {}
            drill_map  = detect_drill_cols(drill_df)   if drill_df  is not None else {}
            delay_map  = detect_delay_cols(delay_df)   if delay_df  is not None else {}

            pptx_buf = generate_report(
                loader_df, loader_map,
                truck_df,  truck_map,
                drill_df,  drill_map,
                delay_df,  delay_map,
                bad_df, week_str, report_date,
            )

        fname = f"PVM_MineStar_Weekly_Report_{report_date.strftime('%d%b%Y')}.pptx"
        st.markdown("""
        <div class="success-banner">
          <h3>✅  Report Ready!</h3>
          <p>Your Capstone Copper-branded PowerPoint is ready to download.</p>
        </div>
        """, unsafe_allow_html=True)
        st.download_button(
            label="⬇️  Download PowerPoint",
            data=pptx_buf,
            file_name=fname,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        st.caption(f"File: {fname}  |  9 slides  |  Capstone Copper template")


if __name__ == "__main__":
    main()
